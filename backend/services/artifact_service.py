#!/usr/bin/env python3
import io
import structlog
from typing import Dict, Any, List, Optional
import uuid
# Imported libraries are assumed to be in requirements.txt (docx, openpyxl, pptx)
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
# Importing utilities to allow for robust PPTX creation (even with mock data)
from pptx.util import Inches

from ..config import get_settings

logger = structlog.get_logger()
settings = get_settings()

class ArtifactService:
    """
    Handles the generation of Project Management artifacts (Excel, Word, PPTX).
    Methods return file bytes.
    """
    def __init__(self):
        pass

    def _generate_excel_risk_register(self, project_name: str, risks: List[Dict[str, Any]]) -> bytes:
        """Generates a simple Excel file (Risk Register) and returns its bytes."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Risk Register"
        
        # Headers
        headers = ["Risk ID", "Description", "Impact", "Probability", "Status", "Mitigation Plan", "Owner"]
        ws.append(headers)

        # Simple data based on mock or LLM summary
        for i, risk in enumerate(risks, 2):
            ws.append([
                i - 1,
                risk.get("description", "Geological uncertainty."),
                risk.get("impact", "High"),
                risk.get("probability", "Medium"),
                "Open",
                risk.get("mitigation", "Increased seismic survey."),
                risk.get("owner", "Geology Team")
            ])

        file_stream = io.BytesIO()
        wb.save(file_stream)
        file_stream.seek(0)
        return file_stream.read()

    def _generate_word_status_report(self, project_name: str, summary_content: str) -> bytes:
        """Generates a simple Word file (Status Report) and returns its bytes."""
        doc = Document()
        
        doc.add_heading(f"Project Status Report: {project_name}", 0)
        doc.add_paragraph(f"Report Generated: {structlog.get_logger().info('Current time').get('timestamp')}")
        
        doc.add_heading("Executive Summary", 1)
        doc.add_paragraph(summary_content)

        doc.add_heading("Spud Prediction Key Insights", 1)
        p = doc.add_paragraph("Based on the latest data for the target Q4 window, the spud predictive model remains highly accurate, maintaining an 85% confidence level for identified sites.")
        
        file_stream = io.BytesIO()
        doc.save(file_stream)
        file_stream.seek(0)
        return file_stream.read()

    def _generate_pptx_executive_pitch(self, project_name: str, key_metrics: Dict[str, Any]) -> bytes:
        """Generates a simple PPTX file (Executive Pitch) and returns its bytes."""
        prs = Presentation()
        
        # Slide 1: Title Slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = f"Executive Pitch: {project_name}"
        slide.placeholders[1].text = "Predicting Spud Activity Success | PMO Analytics"

        # Slide 2: Key Metrics Slide
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = "Key Success Factors"

        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()  
        
        p = tf.add_paragraph()
        p.text = f"Predicted Spud Sites: {key_metrics.get('sites', 'Sites A & C')}"
        p.level = 0
        
        p = tf.add_paragraph()
        p.text = f"Confidence Level: {key_metrics.get('confidence', '85%')}"
        p.level = 1
        
        file_stream = io.BytesIO()
        prs.save(file_stream)
        file_stream.seek(0)
        return file_stream.read()


    async def generate_artifact(self, artifact_type: str, project_name: str, data: Dict[str, Any] = None) -> Optional[bytes]:
        """Routes the generation request to the correct internal method."""
        if data is None:
            data = {}
        
        try:
            if artifact_type == "excel_risk_register":
                risks = data.get("risks", [{"description": data.get("summary", "New risk identified."), "impact": "Medium", "owner": "AI"}])
                return self._generate_excel_risk_register(project_name, risks)
            
            elif artifact_type == "word_status_report":
                summary = data.get("summary", "Project status is Green.")
                return self._generate_word_status_report(project_name, summary)
            
            elif artifact_type == "pptx_executive_pitch":
                key_metrics = data.get("key_metrics", {"confidence": "85%", "sites": "Sites A & C"})
                return self._generate_pptx_executive_pitch(project_name, key_metrics)
            
            else:
                logger.error("Unknown artifact type requested.", type=artifact_type)
                return None
        
        except Exception as e:
            logger.error("Failed to generate artifact.", type=artifact_type, error=str(e))
            raise
